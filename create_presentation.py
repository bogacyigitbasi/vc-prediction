"""
Generate academic presentation for Temporal GNN VC Prediction project.
One slide per approach with detailed examples.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- Colors ---
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x22, 0x22, 0x22)
DARK_GRAY = RGBColor(0x44, 0x44, 0x44)
GRAY = RGBColor(0x77, 0x77, 0x77)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
BLUE = RGBColor(0x2C, 0x5F, 0x8A)
RED = RGBColor(0xB8, 0x33, 0x33)
GREEN = RGBColor(0x2D, 0x8A, 0x4E)
TABLE_HEADER_BG = RGBColor(0x3B, 0x6E, 0x9C)
TABLE_ALT_BG = RGBColor(0xF2, 0xF2, 0xF2)
INVESTOR_COLOR = RGBColor(0x2C, 0x5F, 0x8A)
COMPANY_COLOR = RGBColor(0xE8, 0x8D, 0x3F)
EDGE_COLOR = RGBColor(0xBB, 0xBB, 0xBB)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT
blank_layout = prs.slide_layouts[6]

TOTAL_SLIDES = 25


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.15):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_size=16,
                   color=BLACK, font_name="Calibri", line_spacing=1.4,
                   alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(line, tuple):
            txt, fmt = line
        else:
            txt, fmt = line, {}
        p.text = txt
        p.font.size = Pt(fmt.get("size", font_size))
        p.font.bold = fmt.get("bold", False)
        p.font.color.rgb = fmt.get("color", color)
        p.font.name = font_name
        p.alignment = fmt.get("align", alignment)
        p.space_after = Pt(fmt.get("space_after", 4))
        p.line_spacing = Pt(fmt.get("size", font_size) * line_spacing)
    return txBox


def add_slide_number(slide, num):
    add_text_box(slide, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.3),
                 f"{num}/{TOTAL_SLIDES}", font_size=10, color=GRAY,
                 alignment=PP_ALIGN.RIGHT)


def slide_title(slide, title, subtitle=""):
    add_text_box(slide, Inches(0.7), Inches(0.35), Inches(11), Inches(0.6),
                 title, font_size=28, bold=True, color=BLUE)
    if subtitle:
        add_text_box(slide, Inches(0.7), Inches(0.9), Inches(11), Inches(0.35),
                     subtitle, font_size=14, color=GRAY)
    add_rect(slide, Inches(0.7), Inches(1.25), Inches(2.2), Inches(0.025), BLUE)


def add_table(slide, left, top, width, rows_data, col_widths):
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.38 * n_rows))
    table = table_shape.table
    for i, w in enumerate(col_widths):
        table.columns[i].width = w
    for row_idx, row in enumerate(rows_data):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(13)
                para.font.name = "Calibri"
                para.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
                if row_idx == 0:
                    para.font.bold = True
                    para.font.color.rgb = WHITE
                else:
                    para.font.color.rgb = BLACK
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
    return table_shape


def add_circle(slide, cx, cy, radius, fill_color, label="", font_size=9, font_color=WHITE):
    left = cx - radius
    top = cy - radius
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, radius * 2, radius * 2)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if label:
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.color.rgb = font_color
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
    return shape


def add_line(slide, x1, y1, x2, y2, color=EDGE_COLOR, width=1.5):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


def draw_network(slide, ox, oy, scale=1.0, show_features=False):
    s = scale
    investors = [
        ("Sequoia", Inches(0.0*s), Inches(0.0*s)),
        ("a16z", Inches(0.0*s), Inches(1.4*s)),
        ("YC", Inches(0.0*s), Inches(2.8*s)),
        ("Accel", Inches(0.0*s), Inches(4.2*s)),
    ]
    companies = [
        ("Stripe", Inches(3.5*s), Inches(-0.3*s)),
        ("Airbnb", Inches(3.5*s), Inches(0.9*s)),
        ("Coinbase", Inches(3.5*s), Inches(2.1*s)),
        ("Dropbox", Inches(3.5*s), Inches(3.3*s)),
        ("GitHub", Inches(3.5*s), Inches(4.5*s)),
    ]
    edges = [(0,0),(0,1),(1,0),(1,2),(1,4),(2,1),(2,3),(3,2),(3,4),(3,3)]
    r_inv = Inches(0.38*s)
    r_comp = Inches(0.32*s)
    for ii, ci in edges:
        _, ix, iy = investors[ii]
        _, cx, cy = companies[ci]
        add_line(slide, ox+ix+r_inv, oy+iy, ox+cx-r_comp, oy+cy, EDGE_COLOR, 1.2*s)
    for name, x, y in investors:
        add_circle(slide, ox+x, oy+y, r_inv, INVESTOR_COLOR, name, int(9*s))
    for name, x, y in companies:
        add_circle(slide, ox+x, oy+y, r_comp, COMPANY_COLOR, name, int(8*s))

    if show_features:
        fs = int(7 * s)
        inv_feats = [
            "200 inv, 12yr, diverse",
            "150 inv, 9yr, tech-focus",
            "1000+ inv, 10yr, broad",
            "80 inv, 7yr, enterprise",
        ]
        for i, (name, x, y) in enumerate(investors):
            add_text_box(slide, ox + x - Inches(1.6*s), oy + y - Inches(0.12*s),
                         Inches(1.4*s), Inches(0.25*s),
                         inv_feats[i], font_size=fs, color=GRAY,
                         alignment=PP_ALIGN.RIGHT)

        comp_feats = [
            "$2M, fintech, US",
            "$600K, travel, US",
            "$1M, crypto, US",
            "$1.2M, cloud, US",
            "$100K, dev-tools, US",
        ]
        for i, (name, x, y) in enumerate(companies):
            add_text_box(slide, ox + x + Inches(0.4*s), oy + y - Inches(0.12*s),
                         Inches(1.6*s), Inches(0.25*s),
                         comp_feats[i], font_size=fs, color=GRAY)


# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(blank_layout)
draw_network(slide, Inches(8.5), Inches(2.0), scale=1.1)
add_circle(slide, Inches(8.8), Inches(7.0), Inches(0.15), INVESTOR_COLOR, "", 1)
add_text_box(slide, Inches(9.1), Inches(6.87), Inches(1.2), Inches(0.3), "Investor", font_size=11, color=GRAY)
add_circle(slide, Inches(10.5), Inches(7.0), Inches(0.15), COMPANY_COLOR, "", 1)
add_text_box(slide, Inches(10.8), Inches(6.87), Inches(1.2), Inches(0.3), "Company", font_size=11, color=GRAY)
add_text_box(slide, Inches(1), Inches(1.5), Inches(7), Inches(1.2),
             "Predicting Startup Follow-on Funding", font_size=38, bold=True, color=BLUE)
add_text_box(slide, Inches(1), Inches(2.7), Inches(7), Inches(0.8),
             "Using Temporal Graph Neural Networks on Venture Capital Networks",
             font_size=20, color=DARK_GRAY)
add_rect(slide, Inches(1), Inches(3.6), Inches(3), Inches(0.025), BLUE)
add_multi_text(slide, Inches(1), Inches(4.2), Inches(6), Inches(2.0), [
    ("Dataset: Crunchbase 2015 -- 168K investments, 66K companies, 30K investors", {"size": 14, "color": GRAY}),
    ("Task: Will a startup raise follow-on funding within 24 months?", {"size": 14, "color": GRAY}),
    ("", {"size": 8}),
    ("Tools: PyTorch Geometric  |  XGBoost  |  sentence-transformers  |  Node2Vec", {"size": 13, "color": GRAY}),
])
add_text_box(slide, Inches(1), Inches(6.3), Inches(5), Inches(0.4), "May 2026", font_size=13, color=GRAY)
add_slide_number(slide, 1)


# ============================================================
# SLIDE 2: MOTIVATION
# ============================================================
slide = prs.slides.add_slide(blank_layout)
slide_title(slide, "Why This Problem?", "The venture capital prediction challenge")
add_multi_text(slide, Inches(0.7), Inches(1.6), Inches(5.5), Inches(5.5), [
    ("The Problem", {"size": 20, "bold": True, "color": BLUE, "space_after": 10}),
    ("Most startups fail to raise follow-on funding after Seed/Series A.", {"size": 15}),
    ("Early prediction could help:", {"size": 15, "space_after": 4}),
    ("   - Investors: better portfolio construction", {"size": 14, "color": DARK_GRAY}),
    ("   - Founders: benchmark fundraising likelihood", {"size": 14, "color": DARK_GRAY}),
    ("   - LPs: evaluate fund performance", {"size": 14, "color": DARK_GRAY, "space_after": 16}),
    ("The Gap", {"size": 20, "bold": True, "color": BLUE, "space_after": 10}),
    ("Prior work treats each company as an isolated row in a spreadsheet.", {"size": 15}),
    ("It ignores who invests alongside whom, which investors", {"size": 15}),
    ("back which companies, and how patterns evolve over time.", {"size": 15}),
])
add_multi_text(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(2.5), [
    ("Our Hypothesis", {"size": 20, "bold": True, "color": BLUE, "space_after": 10}),
    ("The VC ecosystem is a network.", {"size": 15, "space_after": 6}),
    ("If Sequoia invests in Stripe, and also in Airbnb,", {"size": 14}),
    ("that tells us something about both companies.", {"size": 14, "space_after": 10}),
    ("Graph Neural Networks can capture this", {"size": 15, "bold": True}),
    ("relational structure that tabular models miss.", {"size": 15}),
])
add_circle(slide, Inches(8.0), Inches(5.0), Inches(0.28), INVESTOR_COLOR, "Seq", 8)
add_circle(slide, Inches(9.8), Inches(4.5), Inches(0.24), COMPANY_COLOR, "Stripe", 7)
add_circle(slide, Inches(9.8), Inches(5.5), Inches(0.24), COMPANY_COLOR, "Airbnb", 7)
add_circle(slide, Inches(11.3), Inches(5.0), Inches(0.24), COMPANY_COLOR, "???", 8, RED)
add_line(slide, Inches(8.28), Inches(4.85), Inches(9.56), Inches(4.5), EDGE_COLOR, 1.3)
add_line(slide, Inches(8.28), Inches(5.15), Inches(9.56), Inches(5.5), EDGE_COLOR, 1.3)
add_line(slide, Inches(8.28), Inches(5.0), Inches(11.06), Inches(5.0), RGBColor(0xDD, 0x55, 0x55), 1.5)
add_text_box(slide, Inches(9.1), Inches(5.95), Inches(3.0), Inches(0.3),
             "predict: will ??? raise again?", font_size=11, bold=True, color=RED)
add_multi_text(slide, Inches(7.0), Inches(6.4), Inches(5.5), Inches(0.8), [
    ("Research Question:", {"size": 15, "bold": True, "color": RED}),
    ("Can GNNs outperform XGBoost by leveraging investor network structure?", {"size": 14}),
])
add_slide_number(slide, 2)


# ============================================================
# SLIDE 3: DATASET
# ============================================================
slide = prs.slides.add_slide(blank_layout)
slide_title(slide, "Dataset", "Crunchbase 2015 relational export")
table_data = [
    ["Component", "Count", "Description"],
    ["Companies", "66,369", "Startups with metadata (sector, country, founded date)"],
    ["Investors", "30,732", "VCs, angels, accelerators"],
    ["Investments", "168,647", "Individual investment records with dates & amounts"],
    ["Funding Rounds", "114,949", "Round-level info (Seed, A, B, ...)"],
    ["Labeled Companies", "23,740", "Had Seed/Series A + observable 24-month window"],
]
add_table(slide, Inches(0.7), Inches(1.6), Inches(8.5), table_data,
          [Inches(2.0), Inches(1.2), Inches(5.3)])
add_multi_text(slide, Inches(0.7), Inches(4.6), Inches(5), Inches(2.5), [
    ("Label Definition", {"size": 17, "bold": True, "color": BLUE, "space_after": 6}),
    ("label = 1  Company raised again within 24 months", {"size": 14}),
    ("label = 0  No follow-on funding in that window", {"size": 14, "space_after": 8}),
    ("Positive rate: 39.1% -- good class balance", {"size": 14, "bold": True}),
])
add_multi_text(slide, Inches(7.0), Inches(4.6), Inches(5.5), Inches(2.5), [
    ("Key Constraint", {"size": 17, "bold": True, "color": RED, "space_after": 6}),
    ("Each company's trigger date + 24 months must fall", {"size": 14}),
    ("before the dataset end (Dec 2015).", {"size": 14, "space_after": 4}),
    ("This ensures every label is fully determined --", {"size": 14}),
    ("no censored observations.", {"size": 14}),
])
add_slide_number(slide, 3)


# ============================================================
# SLIDE 4: GRAPH CONSTRUCTION
# ============================================================
slide = prs.slides.add_slide(blank_layout)
slide_title(slide, "Graph Construction", "Bipartite heterogeneous graph")
add_multi_text(slide, Inches(0.7), Inches(1.6), Inches(5.5), Inches(5.0), [
    ("Two Node Types", {"size": 18, "bold": True, "color": BLUE, "space_after": 8}),
    ("Company nodes  (66,369)  -- startups seeking funding", {"size": 14, "space_after": 2}),
    ("Investor nodes  (30,732)  -- VCs, angels, accelerators", {"size": 14, "space_after": 16}),
    ("Two Edge Types (bidirectional)", {"size": 18, "bold": True, "color": BLUE, "space_after": 8}),
    ("invested_in:      Investor -> Company", {"size": 14, "space_after": 2}),
    ("received_from:  Company -> Investor  (reverse)", {"size": 14, "space_after": 16}),
    ("Why Bipartite?", {"size": 18, "bold": True, "color": BLUE, "space_after": 8}),
    ("Preserves the real structure -- investors and companies", {"size": 14}),
    ("are fundamentally different entity types.", {"size": 14, "space_after": 4}),
    ("A homogeneous projection (company-to-company via shared", {"size": 13, "color": GRAY}),
    ("investors) would lose this distinction.", {"size": 13, "color": GRAY}),
])
draw_network(slide, Inches(7.8), Inches(2.0), scale=1.0, show_features=True)
add_circle(slide, Inches(7.8), Inches(6.7), Inches(0.13), INVESTOR_COLOR, "", 1)
add_text_box(slide, Inches(8.05), Inches(6.58), Inches(1.5), Inches(0.3), "Investor (7 features)", font_size=10, color=GRAY)
add_circle(slide, Inches(9.8), Inches(6.7), Inches(0.13), COMPANY_COLOR, "", 1)
add_text_box(slide, Inches(10.05), Inches(6.58), Inches(1.8), Inches(0.3), "Company (442 features)", font_size=10, color=GRAY)
add_text_box(slide, Inches(7.8), Inches(6.9), Inches(4.5), Inches(0.3),
             "Gray text = sample features on each node", font_size=10, color=GRAY)
add_slide_number(slide, 4)


# ============================================================
# SLIDE 5: TEMPORAL INTEGRITY
# ============================================================
slide = prs.slides.add_slide(blank_layout)
slide_title(slide, "Temporal Integrity", "The most critical design decision")
add_text_box(slide, Inches(0.7), Inches(1.6), Inches(12), Inches(0.4),
             "RULE: Never use future information to predict the past. This is data leakage.",
             font_size=15, bold=True, color=RED)
add_multi_text(slide, Inches(0.7), Inches(2.3), Inches(5.5), Inches(4.8), [
    ("How We Enforce It", {"size": 18, "bold": True, "color": BLUE, "space_after": 8}),
    ("Split by trigger round date (not random):", {"size": 14, "space_after": 6}),
    ("   Train:  trigger before Jan 2012    (12,920 companies)", {"size": 13}),
    ("   Val:     Jan 2012 - Dec 2012          (4,831 companies)", {"size": 13}),
    ("   Test:   Jan 2013 - Dec 2013          (5,989 companies)", {"size": 13, "space_after": 12}),
    ("Each split's graph uses ONLY edges", {"size": 14, "bold": True}),
    ("dated before that split's cutoff.", {"size": 14, "bold": True, "space_after": 12}),
    ("Automated validation checks:", {"size": 14, "space_after": 6}),
    ("   [pass]  No company in multiple splits", {"size": 13, "color": GREEN}),
    ("   [pass]  Train dates < Val dates < Test dates", {"size": 13, "color": GREEN}),
    ("   [pass]  All splits non-empty", {"size": 13, "color": GREEN}),
])
add_multi_text(slide, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.8), [
    ("Timeline", {"size": 16, "bold": True, "color": BLUE, "space_after": 12}),
    ("", {"size": 6}),
    ("  --- 2005 ------- 2012 ------- 2013 ------- 2014 -->", {"size": 14}),
    ("           |               |               |", {"size": 14}),
    ("        Train            Val            Test", {"size": 14, "bold": True, "color": BLUE}),
    ("        cutoff           cutoff          cutoff", {"size": 13, "color": GRAY}),
    ("", {"size": 10}),
    ("  Train graph:  73,509 edges (pre-2012 only)", {"size": 13}),
    ("  Val graph:    91,573 edges (pre-2013 only)", {"size": 13}),
    ("  Test graph:  115,564 edges (pre-2014 only)", {"size": 13}),
    ("", {"size": 10}),
    ("  The graph grows over time -- this is realistic.", {"size": 13, "color": GRAY}),
    ("  A model deployed today wouldn't know about", {"size": 13, "color": GRAY}),
    ("  tomorrow's investments.", {"size": 13, "color": GRAY}),
])
add_slide_number(slide, 5)


# ============================================================
# SLIDE 6: FEATURE ENGINEERING
# ============================================================
slide = prs.slides.add_slide(blank_layout)
slide_title(slide, "Feature Engineering", "What each node knows about itself")
add_multi_text(slide, Inches(0.7), Inches(1.6), Inches(5.8), Inches(2.5), [
    ("Company Features  (58 dimensions)", {"size": 17, "bold": True, "color": BLUE, "space_after": 6}),
    ("Continuous (7):  founded_year, days_since_founded,", {"size": 13}),
    ("   log_total_raised, round_count, investor_count,", {"size": 13}),
    ("   log_max_round_amount, days_since_last_round", {"size": 13, "space_after": 4}),
    ("Binary (1):  is_US", {"size": 13, "space_after": 4}),
    ("Categorical (50):  top-50 industry categories (one-hot)", {"size": 13, "space_after": 4}),
    ("All temporal features computed relative to split cutoff date.", {"size": 12, "color": GRAY}),
    ("Z-score normalized.", {"size": 12, "color": GRAY}),
])
add_multi_text(slide, Inches(0.7), Inches(4.5), Inches(5.8), Inches(2.5), [
    ("Investor Features  (7 dimensions)", {"size": 17, "bold": True, "color": BLUE, "space_after": 6}),
    ("total_investments, portfolio_size, category_diversity,", {"size": 13}),
    ("log_median_amount, years_active, avg_round_stage,", {"size": 13}),
    ("geo_diversity", {"size": 13}),
])
add_multi_text(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.0), [
    ("Text Embeddings  (+384 dimensions)", {"size": 17, "bold": True, "color": BLUE, "space_after": 8}),
    ("Added in V2 to enrich company features.", {"size": 14, "space_after": 6}),
    ("Model: all-MiniLM-L6-v2  (sentence-transformers)", {"size": 13, "space_after": 4}),
    ("Input: company category_list text", {"size": 13}),
    ("   e.g. 'Machine Learning, AI, Analytics'", {"size": 12, "color": GRAY, "space_after": 4}),
    ("Output: 384-dimensional dense vector", {"size": 13, "space_after": 8}),
    ("Why? Captures semantic similarity:", {"size": 14, "bold": True, "space_after": 4}),
    ("   'AI, Machine Learning'  and", {"size": 13, "color": GRAY}),
    ("   'Artificial Intelligence, Analytics'", {"size": 13, "color": GRAY}),
    ("   get similar vectors even with different words.", {"size": 13, "color": GRAY, "space_after": 10}),
    ("Final company input: 58 + 384 = 442 dimensions", {"size": 14, "bold": True}),
])
add_slide_number(slide, 6)


# ============================================================
# SLIDE 7: TOOLS & STACK
# ============================================================
slide = prs.slides.add_slide(blank_layout)
slide_title(slide, "Tools & Technology Stack", "What we used and why")
tools_all = [
    ("PyTorch Geometric (PyG)", "Graph neural network library. Provides GNN layers (SAGEConv, GATConv, GraphConv),\nHeteroData container for bipartite graphs, and to_hetero() wrapper."),
    ("XGBoost", "Gradient-boosted decision trees. Gold standard for tabular data.\nOur baseline -- if GNNs can't beat it, graph structure isn't helping."),
    ("sentence-transformers", "Pre-trained language model (all-MiniLM-L6-v2). Encodes text categories into\n384-dim semantic vectors. No fine-tuning needed -- used off-the-shelf."),
    ("Node2Vec", "Unsupervised graph embedding method. Learns 64-dim vectors from random walks.\nUsed to validate graph structure encodes real patterns (not in prediction models)."),
    ("scikit-learn", "Metrics (F1, AUC, precision, recall), preprocessing (normalization),\nand logistic regression for link prediction task."),
    ("Google Colab (T4 GPU)", "Training environment. Data prepared locally, uploaded as zip to Google Drive,\nmodels trained on free GPU. All experiments use seed=42 for reproducibility."),
]
for i, (name, desc) in enumerate(tools_all):
    y = Inches(1.55) + Inches(i * 0.95)
    add_text_box(slide, Inches(0.7), y, Inches(3.2), Inches(0.3), name, font_size=14, bold=True, color=BLUE)
    add_text_box(slide, Inches(4.0), y, Inches(9.0), Inches(0.85), desc, font_size=12, color=DARK_GRAY)
    if i < len(tools_all) - 1:
        add_rect(slide, Inches(0.7), y + Inches(0.85), Inches(12.0), Inches(0.008), LIGHT_GRAY)
add_slide_number(slide, 7)


# ============================================================
# SLIDE 8: EVALUATION METRICS
# ============================================================
slide = prs.slides.add_slide(blank_layout)
slide_title(slide, "Evaluation Metrics", "How we measure model performance")
add_multi_text(slide, Inches(0.7), Inches(1.6), Inches(5.8), Inches(5.5), [
    ("Precision", {"size": 20, "bold": True, "color": BLUE, "space_after": 6}),
    ("When the model says 'this company will raise', how often is it right?", {"size": 14, "space_after": 4}),
    ("   Precision = True Positives / (True Positives + False Positives)", {"size": 13, "color": GRAY}),
    ("   If precision = 0.55, then 45% of 'yes' predictions are wrong.", {"size": 13, "color": GRAY, "space_after": 16}),
    ("Recall (Sensitivity)", {"size": 20, "bold": True, "color": BLUE, "space_after": 6}),
    ("Of all companies that actually raised, how many did we catch?", {"size": 14, "space_after": 4}),
    ("   Recall = True Positives / (True Positives + False Negatives)", {"size": 13, "color": GRAY}),
    ("   If recall = 0.75, we missed 25% of successful companies.", {"size": 13, "color": GRAY, "space_after": 16}),
    ("F1-Score", {"size": 20, "bold": True, "color": BLUE, "space_after": 6}),
    ("Harmonic mean of precision and recall. Balances both.", {"size": 14, "space_after": 4}),
    ("   F1 = 2 * (Precision * Recall) / (Precision + Recall)", {"size": 13, "color": GRAY}),
    ("   High F1 means good at catching positives without too many false alarms.", {"size": 13, "color": GRAY}),
])
add_multi_text(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.5), [
    ("ROC-AUC", {"size": 20, "bold": True, "color": BLUE, "space_after": 6}),
    ("Area Under the ROC Curve. Measures ranking quality.", {"size": 14, "space_after": 4}),
    ("   AUC = 1.0  Model perfectly separates positive/negative", {"size": 13, "color": GRAY}),
    ("   AUC = 0.5  Model is guessing randomly", {"size": 13, "color": GRAY, "space_after": 4}),
    ("   'How well does the model rank companies that raised", {"size": 13, "color": GRAY}),
    ("    above companies that didn't?'", {"size": 13, "color": GRAY, "space_after": 16}),
    ("Why not Accuracy?", {"size": 20, "bold": True, "color": BLUE, "space_after": 6}),
    ("With 39% positive rate, a model that always says 'no'", {"size": 14}),
    ("gets 61% accuracy for free. Accuracy is misleading", {"size": 14}),
    ("with class imbalance -- it doesn't tell you whether the", {"size": 14}),
    ("model actually learned anything useful.", {"size": 14, "space_after": 16}),
    ("Our primary metrics: F1 and ROC-AUC", {"size": 16, "bold": True, "color": BLUE, "space_after": 6}),
    ("F1 balances precision/recall.", {"size": 14}),
    ("AUC measures overall ranking ability.", {"size": 14}),
])
add_slide_number(slide, 8)


# ============================================================
# SLIDE 9: HOW XGBOOST WORKS
# ============================================================
slide = prs.slides.add_slide(blank_layout)
slide_title(slide, "Approach 1: XGBoost", "Gradient-boosted decision trees (our baseline)")

add_multi_text(slide, Inches(0.7), Inches(1.6), Inches(5.5), Inches(5.5), [
    ("What it does", {"size": 18, "bold": True, "color": BLUE, "space_after": 8}),
    ("Builds 500 small decision trees, one after another.", {"size": 14}),
    ("Each new tree focuses on the mistakes of the previous ones.", {"size": 14, "space_after": 10}),
    ("Step-by-step example:", {"size": 16, "bold": True, "color": BLUE, "space_after": 6}),
    ("   Tree #1 asks:", {"size": 14, "bold": True}),
    ("   Is total_raised > $1M?", {"size": 13}),
    ("     YES -> Is investor_count > 3?", {"size": 13, "color": GRAY}),
    ("       YES -> probably raises again (label=1)", {"size": 13, "color": GREEN}),
    ("       NO  -> maybe not (label=0)", {"size": 13, "color": RED}),
    ("     NO  -> probably doesn't (label=0)", {"size": 13, "color": RED, "space_after": 8}),
    ("   Tree #1 gets 200 companies wrong.", {"size": 13}),
    ("   Tree #2 focuses only on those 200 mistakes.", {"size": 13}),
    ("   Tree #3 focuses on remaining errors.", {"size": 13}),
    ("   ...500 trees later, errors are minimized.", {"size": 13}),
])

add_multi_text(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(2.5), [
    ("What it sees", {"size": 18, "bold": True, "color": BLUE, "space_after": 8}),
    ("ONLY the company's own features (58 numbers):", {"size": 14}),
    ("   raised=$500K, investors=2, sector=fintech, US=yes ...", {"size": 13, "color": GRAY, "space_after": 8}),
    ("It has NO idea who the investors are,", {"size": 14, "bold": True, "color": RED}),
    ("what else they've invested in, or whether", {"size": 14, "bold": True, "color": RED}),
    ("their other companies succeeded.", {"size": 14, "bold": True, "color": RED}),
])

add_multi_text(slide, Inches(7.0), Inches(4.5), Inches(5.5), Inches(2.5), [
    ("Why it's the baseline", {"size": 18, "bold": True, "color": BLUE, "space_after": 8}),
    ("XGBoost is the gold standard for tabular data.", {"size": 14}),
    ("If our GNNs can't beat this, then graph structure", {"size": 14}),
    ("is not adding value. This is our scientific control.", {"size": 14, "space_after": 8}),
    ("Result: F1 = 0.558, AUC = 0.711", {"size": 15, "bold": True}),
])
add_slide_number(slide, 9)


# ============================================================
# SLIDE 10: HOW GCN WORKS
# ============================================================
slide = prs.slides.add_slide(blank_layout)
slide_title(slide, "Approach 2: GCN (Graph Convolutional Network)", "Symmetric neighborhood averaging")

add_multi_text(slide, Inches(0.7), Inches(1.6), Inches(5.8), Inches(5.5), [
    ("What it does", {"size": 18, "bold": True, "color": BLUE, "space_after": 8}),
    ("Each node updates its representation by averaging", {"size": 14}),
    ("ALL of its neighbors' features, weighted by degree.", {"size": 14, "space_after": 10}),
    ("Step-by-step example:", {"size": 16, "bold": True, "color": BLUE, "space_after": 6}),
    ("", {"size": 4}),
    ("   Company X has 3 investors: Sequoia, YC, an angel", {"size": 14, "space_after": 6}),
    ("   Layer 1:", {"size": 14, "bold": True}),
    ("   new_X = Average(Sequoia_features, YC_features, angel_features)", {"size": 13, "color": GRAY}),
    ("   All 3 investors weighted equally by their degree.", {"size": 13, "color": GRAY, "space_after": 6}),
    ("   Layer 2:", {"size": 14, "bold": True}),
    ("   Sequoia was also updated using ITS portfolio companies.", {"size": 13, "color": GRAY}),
    ("   So now X indirectly 'knows' about Sequoia's other bets.", {"size": 13, "color": GRAY, "space_after": 10}),
    ("   Final: updated_X -> MLP classifier -> P(follow-on)", {"size": 14, "bold": True}),
])

add_multi_text(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(2.5), [
    ("The key idea", {"size": 18, "bold": True, "color": BLUE, "space_after": 8}),
    ("h_v = sigma( D^(-1/2) * A * D^(-1/2) * H * W )", {"size": 13, "color": GRAY, "space_after": 6}),
    ("In plain English:", {"size": 14}),
    ("'My new representation is the weighted average", {"size": 14}),
    (" of my neighbors, passed through a neural network.'", {"size": 14, "space_after": 6}),
    ("Every neighbor gets equal importance.", {"size": 14, "bold": True}),
])

add_multi_text(slide, Inches(7.0), Inches(4.5), Inches(5.5), Inches(2.5), [
    ("Problem on our graph", {"size": 18, "bold": True, "color": RED, "space_after": 8}),
    ("GCN needs self-loops (a node connects to itself).", {"size": 14}),
    ("Self-loops require a square adjacency matrix.", {"size": 14}),
    ("Our bipartite graph is NOT square (66K x 30K).", {"size": 14, "space_after": 6}),
    ("We used GraphConv instead, but training was unstable --", {"size": 13}),
    ("loss spiked to 9.0. GCN failed to converge.", {"size": 13, "bold": True, "color": RED}),
])
add_slide_number(slide, 10)


# ============================================================
# SLIDE 11: HOW GRAPHSAGE WORKS
# ============================================================
slide = prs.slides.add_slide(blank_layout)
slide_title(slide, "Approach 3: GraphSAGE", "Sample-and-aggregate (our best GNN)")

add_multi_text(slide, Inches(0.7), Inches(1.6), Inches(5.8), Inches(5.5), [
    ("What it does", {"size": 18, "bold": True, "color": BLUE, "space_after": 8}),
    ("Instead of averaging ALL neighbors equally, GraphSAGE:", {"size": 14}),
    ("  1. Samples a fixed number of neighbors", {"size": 14}),
    ("  2. Computes the mean of their features", {"size": 14}),
    ("  3. Concatenates with the node's own features", {"size": 14}),
    ("  4. Applies a learned transformation (W matrix)", {"size": 14, "space_after": 10}),
    ("Step-by-step example:", {"size": 16, "bold": True, "color": BLUE, "space_after": 6}),
    ("", {"size": 4}),
    ("   Company X has 3 investors: Sequoia, YC, angel", {"size": 14, "space_after": 6}),
    ("   Layer 1:", {"size": 14, "bold": True}),
    ("   neighbor_mean = Mean(Sequoia_feat, YC_feat, angel_feat)", {"size": 13, "color": GRAY}),
    ("   new_X = W * CONCAT(X_own_feat, neighbor_mean)", {"size": 13, "color": GRAY, "space_after": 6}),
    ("   Layer 2:", {"size": 14, "bold": True}),
    ("   Each investor was updated from THEIR portfolio companies.", {"size": 13, "color": GRAY}),
    ("   Sequoia now carries info from Stripe, Airbnb, etc.", {"size": 13, "color": GRAY}),
    ("   X absorbs this 2nd-hop information.", {"size": 13, "color": GRAY}),
])

add_multi_text(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(2.5), [
    ("Why it works best on our graph", {"size": 18, "bold": True, "color": GREEN, "space_after": 8}),
    ("No self-loops needed -- works on bipartite graphs", {"size": 14}),
    ("Handles power-law: some investors have 100+ companies,", {"size": 14}),
    ("   most have <5. Sampling handles this naturally.", {"size": 14, "space_after": 6}),
    ("Inductive: learns a FUNCTION, not node-specific", {"size": 14, "bold": True}),
    ("   embeddings. Can predict on new companies.", {"size": 14}),
])

add_multi_text(slide, Inches(7.0), Inches(4.5), Inches(5.5), Inches(2.5), [
    ("The difference from GCN", {"size": 18, "bold": True, "color": BLUE, "space_after": 8}),
    ("GCN: new_X = WeightedAvg(all neighbors)", {"size": 14}),
    ("SAGE: new_X = W * [X ; Mean(neighbors)]", {"size": 14, "space_after": 6}),
    ("SAGE keeps the node's own features separate,", {"size": 14}),
    ("then combines them. This is more expressive.", {"size": 14, "space_after": 8}),
    ("Result: F1 = 0.628, AUC = 0.741", {"size": 15, "bold": True, "color": GREEN}),
])
add_slide_number(slide, 11)


# ============================================================
# SLIDE 12: HOW GAT WORKS
# ============================================================
slide = prs.slides.add_slide(blank_layout)
slide_title(slide, "Approach 4: GAT (Graph Attention Network)", "Learning which neighbors matter more")

add_multi_text(slide, Inches(0.7), Inches(1.6), Inches(5.8), Inches(5.5), [
    ("What it does", {"size": 18, "bold": True, "color": BLUE, "space_after": 8}),
    ("Instead of treating all neighbors equally,", {"size": 14}),
    ("GAT learns ATTENTION WEIGHTS -- how much each", {"size": 14}),
    ("neighbor should influence the node.", {"size": 14, "space_after": 10}),
    ("Step-by-step example:", {"size": 16, "bold": True, "color": BLUE, "space_after": 6}),
    ("", {"size": 4}),
    ("   Company X has 3 investors:", {"size": 14, "space_after": 4}),
    ("   GAT learns attention scores:", {"size": 14}),
    ("     Sequoia:  attention = 0.6  (very relevant)", {"size": 13, "color": GRAY}),
    ("     YC:       attention = 0.3  (somewhat relevant)", {"size": 13, "color": GRAY}),
    ("     angel:    attention = 0.1  (less relevant)", {"size": 13, "color": GRAY, "space_after": 6}),
    ("   new_X = 0.6*Sequoia + 0.3*YC + 0.1*angel", {"size": 13, "color": GRAY, "space_after": 8}),
    ("   Multi-head: we run 4 parallel attention heads,", {"size": 14}),
    ("   each learning different importance patterns,", {"size": 14}),
    ("   then concatenate the results.", {"size": 14}),
])

add_multi_text(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(2.3), [
    ("The promise", {"size": 18, "bold": True, "color": BLUE, "space_after": 8}),
    ("Not all investors are equally informative.", {"size": 14}),
    ("Sequoia's signal should matter more than a", {"size": 14}),
    ("random angel. GAT can learn this automatically.", {"size": 14, "space_after": 6}),
    ("In theory, this is more powerful than GCN/SAGE", {"size": 14}),
    ("because it's data-driven, not uniform.", {"size": 14}),
])

add_multi_text(slide, Inches(7.0), Inches(4.3), Inches(5.5), Inches(3.0), [
    ("What went wrong on our graph", {"size": 18, "bold": True, "color": RED, "space_after": 8}),
    ("V1: Underfitting. Flat loss, poor predictions.", {"size": 14}),
    ("   Attention needs more data to learn good weights.", {"size": 13, "color": GRAY, "space_after": 6}),
    ("V2: Degeneration. When we added class weights,", {"size": 14}),
    ("   GAT predicted 'positive' for 99% of companies.", {"size": 14}),
    ("   Attention collapsed -- every investor got equal weight.", {"size": 13, "color": GRAY, "space_after": 6}),
    ("Our bipartite graph with ~24K labeled nodes may simply", {"size": 13}),
    ("not have enough signal for attention to learn from.", {"size": 13, "space_after": 6}),
    ("Result: F1 = 0.519, AUC = 0.596", {"size": 15, "bold": True, "color": RED}),
])
add_slide_number(slide, 12)


# ============================================================
# SLIDE 13: HOW ENSEMBLE WORKS
# ============================================================
slide = prs.slides.add_slide(blank_layout)
slide_title(slide, "Approach 5: Ensemble", "Combining XGBoost + GraphSAGE")

add_multi_text(slide, Inches(0.7), Inches(1.6), Inches(5.8), Inches(5.5), [
    ("What it does", {"size": 18, "bold": True, "color": BLUE, "space_after": 8}),
    ("Run both XGBoost and GraphSAGE separately,", {"size": 14}),
    ("then combine their probability predictions:", {"size": 14, "space_after": 6}),
    ("   P_final = w * P_XGBoost + (1-w) * P_GraphSAGE", {"size": 15, "bold": True, "space_after": 6}),
    ("   w is optimized on the validation set (grid search).", {"size": 13, "color": GRAY, "space_after": 12}),
    ("Step-by-step example:", {"size": 16, "bold": True, "color": BLUE, "space_after": 6}),
    ("", {"size": 4}),
    ("   Company X: $500K seed, 2 investors, fintech, US", {"size": 14, "space_after": 6}),
    ("   XGBoost looks at the spreadsheet row:", {"size": 14}),
    ("     '$500K is below average, only 2 investors... risky'", {"size": 13, "color": GRAY}),
    ("     P_XGBoost = 0.45 -> NO", {"size": 13, "color": RED, "space_after": 6}),
    ("   GraphSAGE looks at the network:", {"size": 14}),
    ("     'But Investor A backed 15 companies and 12 raised!", {"size": 13, "color": GRAY}),
    ("      This company is in a strong syndicate.'", {"size": 13, "color": GRAY}),
    ("     P_GraphSAGE = 0.72 -> YES", {"size": 13, "color": GREEN, "space_after": 6}),
    ("   Ensemble (w=0.4):", {"size": 14}),
    ("     P_final = 0.4*0.45 + 0.6*0.72 = 0.61 -> YES", {"size": 14, "bold": True}),
])

add_multi_text(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(2.5), [
    ("Why it works", {"size": 18, "bold": True, "color": BLUE, "space_after": 8}),
    ("XGBoost and GraphSAGE make different kinds of errors.", {"size": 14, "space_after": 6}),
    ("XGBoost is good at:", {"size": 14, "bold": True}),
    ("   Non-linear tabular patterns", {"size": 13, "color": GRAY}),
    ("   'High funding + few investors = concentrated bet'", {"size": 13, "color": GRAY, "space_after": 6}),
    ("GraphSAGE is good at:", {"size": 14, "bold": True}),
    ("   Relational patterns invisible to tabular models", {"size": 13, "color": GRAY}),
    ("   'Backed by investors whose portfolios succeed'", {"size": 13, "color": GRAY}),
])

add_multi_text(slide, Inches(7.0), Inches(4.5), Inches(5.5), Inches(2.5), [
    ("Averaging reduces variance and captures both.", {"size": 14, "space_after": 12}),
    ("This is our best model.", {"size": 16, "bold": True, "color": GREEN, "space_after": 8}),
    ("Result: F1 = 0.638, AUC = 0.745", {"size": 15, "bold": True, "color": GREEN}),
    ("   +14.3% F1 over XGBoost alone", {"size": 14, "color": DARK_GRAY}),
    ("   Catches 75% of companies that raise (vs 56%)", {"size": 14, "color": DARK_GRAY}),
])
add_slide_number(slide, 13)


# ============================================================
# SLIDE 14: HOW NODE2VEC WORKS
# ============================================================
slide = prs.slides.add_slide(blank_layout)
slide_title(slide, "Approach 6: Node2Vec", "Unsupervised graph embeddings (validation only -- not in prediction models)")

add_multi_text(slide, Inches(0.7), Inches(1.6), Inches(5.8), Inches(5.5), [
    ("What it does", {"size": 18, "bold": True, "color": BLUE, "space_after": 8}),
    ("Learns a 64-dimensional vector for each node by doing", {"size": 14}),
    ("random walks on the graph (like Word2Vec for graphs).", {"size": 14, "space_after": 10}),
    ("Step-by-step example:", {"size": 16, "bold": True, "color": BLUE, "space_after": 6}),
    ("", {"size": 4}),
    ("   Random walk starting from Stripe:", {"size": 14, "space_after": 4}),
    ("   Stripe -> Sequoia -> Airbnb -> YC -> Dropbox -> ...", {"size": 14, "color": GRAY, "space_after": 4}),
    ("   (hop from company to investor to company to investor...)", {"size": 13, "color": GRAY, "space_after": 6}),
    ("   Do this 10 times per node, 20 steps each walk.", {"size": 14, "space_after": 6}),
    ("   Treat each walk as a 'sentence', nodes as 'words'.", {"size": 14}),
    ("   Train Word2Vec: nodes that co-occur in walks", {"size": 14}),
    ("   get similar 64-dim vectors.", {"size": 14, "space_after": 10}),
    ("   Result: Stripe's vector is close to Weebly, Okta, Fivestars", {"size": 14}),
    ("   (all fintech/payments companies with shared investors).", {"size": 14}),
])

add_multi_text(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(2.5), [
    ("Why we use it", {"size": 18, "bold": True, "color": BLUE, "space_after": 8}),
    ("NOT used in V1, V2, or V3 prediction models.", {"size": 14, "bold": True, "color": RED, "space_after": 6}),
    ("It's evidence for our paper. A reviewer would ask:", {"size": 14}),
    ("'How do you know the graph has real structure?'", {"size": 14, "space_after": 6}),
    ("Node2Vec proves it: similar companies cluster together.", {"size": 14}),
    ("This validates the GNN approach before we build it.", {"size": 14}),
])

add_multi_text(slide, Inches(7.0), Inches(4.5), Inches(5.5), Inches(2.5), [
    ("Nearest neighbors found", {"size": 18, "bold": True, "color": BLUE, "space_after": 8}),
    ("   Stripe  -> Weebly, Okta, Fivestars", {"size": 14}),
    ("   Airbnb  -> Blippy, Mixpanel, YC", {"size": 14}),
    ("   Uber     -> Swipely, Gobble, Embedly", {"size": 14, "space_after": 8}),
    ("Companies with shared investors cluster together.", {"size": 14, "bold": True, "color": GREEN}),
    ("The graph structure is real and meaningful.", {"size": 14, "bold": True, "color": GREEN}),
])
add_slide_number(slide, 14)


# ============================================================
# SLIDE 15: HOW SENTENCE-TRANSFORMERS / LINK PREDICTION WORKS
# ============================================================
slide = prs.slides.add_slide(blank_layout)
slide_title(slide, "Approach 7: Sentence-Transformers for Link Prediction", "Predicting WHO invests in WHOM using text embeddings")

add_multi_text(slide, Inches(0.7), Inches(1.6), Inches(5.8), Inches(5.5), [
    ("What it does", {"size": 18, "bold": True, "color": BLUE, "space_after": 8}),
    ("Uses a pre-trained language model to encode company", {"size": 14}),
    ("categories into vectors, then predicts investment links.", {"size": 14, "space_after": 10}),
    ("Step-by-step example:", {"size": 16, "bold": True, "color": BLUE, "space_after": 6}),
    ("", {"size": 4}),
    ("   1. Encode each company's category text:", {"size": 14}),
    ("      'Machine Learning, AI' -> [0.23, -0.41, 0.67, ...]  (384 dims)", {"size": 12, "color": GRAY}),
    ("      'Enterprise SaaS'      -> [0.19, -0.38, 0.71, ...]  (similar!)", {"size": 12, "color": GRAY}),
    ("      'Biotech, Healthcare'  -> [-0.55, 0.82, 0.11, ...]  (different)", {"size": 12, "color": GRAY, "space_after": 6}),
    ("   2. Represent each investor as the mean of their portfolio:", {"size": 14}),
    ("      Sequoia_vec = Average(Stripe_vec, Airbnb_vec, ...)", {"size": 12, "color": GRAY}),
    ("      This captures the investor's 'thesis' -- what they invest in.", {"size": 12, "color": GRAY, "space_after": 6}),
    ("   3. For each investor-company pair, create feature:", {"size": 14}),
    ("      [investor_384d | company_384d | cosine_similarity] = 769d", {"size": 12, "color": GRAY, "space_after": 6}),
    ("   4. Train logistic regression: does this edge exist?", {"size": 14}),
])

add_multi_text(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(2.3), [
    ("A different question than node classification", {"size": 18, "bold": True, "color": BLUE, "space_after": 8}),
    ("Node classification:  'Will Company X raise again?'", {"size": 14}),
    ("Link prediction:       'Will Investor A invest in Company X?'", {"size": 14, "space_after": 8}),
    ("This is about predicting the graph structure itself,", {"size": 14}),
    ("not the outcome of nodes in the graph.", {"size": 14}),
])

add_multi_text(slide, Inches(7.0), Inches(4.3), Inches(5.5), Inches(3.0), [
    ("Practical application", {"size": 18, "bold": True, "color": BLUE, "space_after": 8}),
    ("Deal sourcing tool:", {"size": 14, "bold": True}),
    ("Find companies that match an investor's historical", {"size": 14}),
    ("preferences but haven't appeared in their pipeline.", {"size": 14, "space_after": 8}),
    ("Example: if a VC has only invested in ML/AI companies,", {"size": 13, "color": GRAY}),
    ("their embedding will be close to new ML companies they", {"size": 13, "color": GRAY}),
    ("haven't seen yet -> recommendation engine.", {"size": 13, "color": GRAY}),
])
add_slide_number(slide, 15)


# ============================================================
# SLIDE 16: TRAINING IMPROVEMENTS
# ============================================================
slide = prs.slides.add_slide(blank_layout)
slide_title(slide, "Iterative Improvements", "Three rounds of training: V1 -> V2 -> V3")
versions = [
    ("V1: Baseline GNNs", [
        "BCELoss with sigmoid",
        "Adam optimizer (lr=0.005)",
        "Early stopping on val F1",
        "Patience: 20 epochs",
    ]),
    ("V2: Stabilization", [
        "BCEWithLogitsLoss + class weights",
        "Gradient clipping (max_norm=1.0)",
        "Per-model learning rates",
        "LR scheduler (ReduceLROnPlateau)",
        "Added text embeddings (+384d)",
    ]),
    ("V3: Optimization", [
        "3-layer deep GraphSAGE",
        "Residual connections + BatchNorm",
        "Threshold optimization (PR curve)",
        "XGBoost + GraphSAGE ensemble",
        "Grid search for ensemble weights",
    ]),
]
for i, (title, items) in enumerate(versions):
    x = Inches(0.5) + Inches(i * 4.2)
    add_rect(slide, x, Inches(1.6), Inches(3.9), Inches(4.2), WHITE, LIGHT_GRAY)
    add_text_box(slide, x + Inches(0.2), Inches(1.75), Inches(3.5), Inches(0.35),
                 title, font_size=16, bold=True, color=BLUE)
    add_rect(slide, x + Inches(0.2), Inches(2.15), Inches(1.8), Inches(0.015), LIGHT_GRAY)
    items_text = "\n".join(f"  - {item}" for item in items)
    add_text_box(slide, x + Inches(0.2), Inches(2.35), Inches(3.5), Inches(3.2),
                 items_text, font_size=13, color=DARK_GRAY)
add_text_box(slide, Inches(0.7), Inches(6.1), Inches(12.0), Inches(0.35),
             "Key: GCN and GAT degenerate under class weighting (predict all positive), but GraphSAGE stays stable.",
             font_size=13, color=RED)
add_text_box(slide, Inches(0.7), Inches(6.5), Inches(12.0), Inches(0.35),
             "Deep GraphSAGE (3 layers) slightly underperforms 2 layers -- over-smoothing on small-diameter bipartite graphs.",
             font_size=13, color=GRAY)
add_slide_number(slide, 16)


# ============================================================
# SLIDE 17: RESULTS TABLE
# ============================================================
slide = prs.slides.add_slide(blank_layout)
slide_title(slide, "Results: Full Comparison", "Test set performance across all models")
results_data = [
    ["Model", "F1", "ROC-AUC", "Precision", "Recall", "Notes"],
    ["XGBoost (baseline)", "0.558", "0.711", "0.557", "0.558", "No graph structure"],
    ["GCN v1", "--", "--", "--", "--", "Failed to converge"],
    ["GAT v1", "0.519", "0.596", "0.430", "0.654", "Underfitting"],
    ["GraphSAGE v1", "0.625", "0.710", "0.542", "0.739", "+12% F1 vs. XGBoost"],
    ["GraphSAGE v2", "0.628", "0.741", "0.543", "0.744", "Best single GNN"],
    ["Deep GraphSAGE v3", "0.619", "0.738", "0.536", "0.732", "Over-smoothing"],
    ["Ensemble (XGB+v2)", "0.638", "0.745", "0.557", "0.745", "Best overall"],
]
table_s = add_table(slide, Inches(0.7), Inches(1.6), Inches(12.0), results_data,
          [Inches(2.5), Inches(0.8), Inches(1.0), Inches(1.0), Inches(0.8), Inches(5.9)])
table = table_s.table
for col_idx in range(6):
    cell = table.cell(7, col_idx)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0xD5, 0xF5, 0xE3)
add_multi_text(slide, Inches(0.7), Inches(5.5), Inches(12.0), Inches(1.5), [
    ("Ensemble formula:", {"size": 15, "bold": True, "color": BLUE, "space_after": 4}),
    ("P_ensemble  =  w * P_XGBoost  +  (1-w) * P_GraphSAGE      (w optimized on validation set)", {"size": 14}),
    ("", {"size": 6}),
    ("XGBoost captures tabular interactions (e.g., 'high funding + few investors = concentrated bet')", {"size": 13, "color": GRAY}),
    ("GraphSAGE captures relational patterns (e.g., 'backed by investors whose portfolios succeed')", {"size": 13, "color": GRAY}),
])
add_slide_number(slide, 17)


# ============================================================
# SLIDE 18: PERFORMANCE JOURNEY
# ============================================================
slide = prs.slides.add_slide(blank_layout)
slide_title(slide, "Performance Journey", "From baseline to best model: +14.3% F1")
steps = [
    ("XGBoost\nbaseline", "F1 = 0.558\nAUC = 0.711", "floor"),
    ("+ Graph\nstructure", "F1 = 0.625\nAUC = 0.710", "+12.0% F1"),
    ("+ Text\n+ weights", "F1 = 0.628\nAUC = 0.741", "+0.5% F1"),
    ("+ Ensemble\n(XGB+SAGE)", "F1 = 0.638\nAUC = 0.745", "+1.6% F1"),
]
for i, (label, metric, delta) in enumerate(steps):
    x = Inches(0.7) + Inches(i * 3.1)
    add_rect(slide, x, Inches(2.5), Inches(2.5), Inches(3.0), WHITE, LIGHT_GRAY)
    add_text_box(slide, x + Inches(0.15), Inches(2.65), Inches(2.2), Inches(0.6),
                 label, font_size=14, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)
    add_rect(slide, x + Inches(0.3), Inches(3.35), Inches(1.9), Inches(0.015), LIGHT_GRAY)
    add_text_box(slide, x + Inches(0.15), Inches(3.55), Inches(2.2), Inches(0.7),
                 metric, font_size=16, bold=True, color=BLACK, alignment=PP_ALIGN.CENTER)
    delta_color = GRAY if i == 0 else GREEN
    add_text_box(slide, x + Inches(0.15), Inches(4.6), Inches(2.2), Inches(0.4),
                 delta, font_size=13, color=delta_color, alignment=PP_ALIGN.CENTER)
    if i < len(steps) - 1:
        add_text_box(slide, x + Inches(2.5), Inches(3.5), Inches(0.6), Inches(0.4),
                     "->", font_size=22, color=GRAY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(0.7), Inches(6.0), Inches(12.0), Inches(0.5),
             "Total improvement:  F1 +14.3%  |  AUC +4.8%  |  Recall 56% -> 75%  (catches 3 out of 4 successful startups)",
             font_size=15, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 18)


# ============================================================
# SLIDE 19: ABLATION STUDIES
# ============================================================
slide = prs.slides.add_slide(blank_layout)
slide_title(slide, "Ablation Studies", "Proving what actually matters")
abl_data = [
    ["Setting", "F1", "AUC", "What's Removed"],
    ["Full model (SAGEv2)", "0.628", "0.741", "Nothing (control)"],
    ["Topology-only", "0.572", "0.544", "All node features (random features)"],
    ["No temporal filter", "0.620", "0.730", "Temporal edge filtering"],
    ["No graph (XGBoost)", "0.558", "0.711", "All graph structure"],
]
add_table(slide, Inches(0.7), Inches(1.6), Inches(12.0), abl_data,
          [Inches(3.0), Inches(1.0), Inches(1.0), Inches(7.0)])
add_multi_text(slide, Inches(0.7), Inches(3.8), Inches(5.5), Inches(3.5), [
    ("Finding 1: Features dominate", {"size": 17, "bold": True, "color": BLUE, "space_after": 6}),
    ("Removing features collapses AUC from 0.74 to 0.54", {"size": 13}),
    ("(barely above 0.50 random chance).", {"size": 13, "space_after": 4}),
    ("Node features carry ~20pp of the signal.", {"size": 13, "bold": True, "space_after": 14}),
    ("Finding 2: Graph adds complementary signal", {"size": 17, "bold": True, "color": BLUE, "space_after": 6}),
    ("Graph structure provides +3pp AUC over XGBoost.", {"size": 13}),
    ("Smaller than features, but consistent and valuable.", {"size": 13}),
])
add_multi_text(slide, Inches(7.0), Inches(3.8), Inches(5.5), Inches(3.5), [
    ("Finding 3: Temporal filtering is sound", {"size": 17, "bold": True, "color": BLUE, "space_after": 6}),
    ("Leaking future edges doesn't inflate metrics.", {"size": 13, "space_after": 4}),
    ("VC network structure is stable over 2012-2014.", {"size": 13}),
    ("The methodology works.", {"size": 13, "space_after": 14}),
    ("Signal Decomposition", {"size": 17, "bold": True, "color": BLUE, "space_after": 6}),
    ("   Features:           ~20pp AUC (primary)", {"size": 14}),
    ("   Graph topology:  ~3pp AUC (complementary)", {"size": 14}),
    ("   Both together:    0.741 AUC (best)", {"size": 14, "bold": True}),
])
add_slide_number(slide, 19)


# ============================================================
# SLIDE 20: WHY GRAPHSAGE WINS
# ============================================================
slide = prs.slides.add_slide(blank_layout)
slide_title(slide, "Discussion: Why GraphSAGE?", "Architecture analysis")
reasons = [
    ("Bipartite-friendly", "Mean aggregation works naturally\non investor -> company edges.\nNo self-loops needed\n(unlike GCN)."),
    ("Power-law robust", "VC networks are scale-free:\na few investors back 100s of\ncompanies, most back <5.\nSAGE handles this gracefully."),
    ("Inductive", "Generalizes to unseen nodes\nat inference time. Critical for\nproduction: new startups\nappear every day."),
    ("Stable training", "Handles class weighting,\nvarying learning rates, and\nbipartite structure without\ndegeneration."),
]
for i, (title, desc) in enumerate(reasons):
    x = Inches(0.4) + Inches(i * 3.15)
    add_rect(slide, x, Inches(1.6), Inches(2.9), Inches(3.2), WHITE, LIGHT_GRAY)
    add_text_box(slide, x + Inches(0.15), Inches(1.75), Inches(2.6), Inches(0.35),
                 title, font_size=15, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)
    add_rect(slide, x + Inches(0.3), Inches(2.15), Inches(2.3), Inches(0.015), LIGHT_GRAY)
    add_text_box(slide, x + Inches(0.15), Inches(2.35), Inches(2.6), Inches(2.2),
                 desc, font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_multi_text(slide, Inches(0.7), Inches(5.2), Inches(12.0), Inches(2.0), [
    ("Over-smoothing effect:", {"size": 15, "bold": True, "color": BLUE, "space_after": 4}),
    ("Deep GraphSAGE (3 layers + residual + batchnorm) slightly underperforms 2-layer version.", {"size": 13}),
    ("Bipartite graphs have small diameter -- after 3 hops, all company nodes converge to similar representations.", {"size": 13}),
    ("Lesson: deeper is not always better for GNNs. 2 layers is the sweet spot for this graph topology.", {"size": 13, "color": GRAY}),
])
add_slide_number(slide, 20)


# ============================================================
# SLIDE 21: DEMO
# ============================================================
slide = prs.slides.add_slide(blank_layout)
slide_title(slide, "Demo: How It Works in Practice", "From raw data to prediction")
steps_demo = [
    ("1. Data Prep\n(Local)", "Load Crunchbase CSVs\nGenerate labels\nTemporal split\nBuild edge lists\nCompute features"),
    ("2. Upload\n(Colab)", "Zip processed data\nMount Google Drive\nExtract on GPU server\nBuild PyG HeteroData"),
    ("3. Train\n(GPU)", "XGBoost baseline\nGraphSAGE (2 layers)\nEarly stopping\nLR scheduling\nGradient clipping"),
    ("4. Evaluate\n(Colab)", "Test set predictions\nF1, AUC, precision, recall\nAblation studies\nEnsemble optimization"),
]
for i, (title, desc) in enumerate(steps_demo):
    x = Inches(0.4) + Inches(i * 3.15)
    add_rect(slide, x, Inches(1.6), Inches(2.9), Inches(3.0), WHITE, BLUE)
    add_text_box(slide, x + Inches(0.15), Inches(1.75), Inches(2.6), Inches(0.6),
                 title, font_size=14, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)
    add_rect(slide, x + Inches(0.3), Inches(2.35), Inches(2.3), Inches(0.015), LIGHT_GRAY)
    add_text_box(slide, x + Inches(0.15), Inches(2.55), Inches(2.6), Inches(1.8),
                 desc, font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
    if i < len(steps_demo) - 1:
        add_text_box(slide, x + Inches(2.9), Inches(2.6), Inches(0.3), Inches(0.4),
                     "->", font_size=20, color=BLUE, alignment=PP_ALIGN.CENTER)
add_multi_text(slide, Inches(0.7), Inches(5.0), Inches(12.0), Inches(2.2), [
    ("Example Prediction (test set)", {"size": 15, "bold": True, "color": BLUE, "space_after": 8}),
    ("Input:  Company X -- raised $500K Seed, 2 investors (1 active VC, 1 angel), fintech sector, US-based", {"size": 13, "space_after": 6}),
    ("XGBoost says:     P(follow-on) = 0.45  ->  NO   (only sees company features)", {"size": 13, "color": GRAY}),
    ("GraphSAGE says:  P(follow-on) = 0.72  ->  YES  (sees that the VC's other companies mostly raised)", {"size": 13, "color": GRAY}),
    ("Ensemble says:    P(follow-on) = 0.61  ->  YES  (weighted average, more confident)", {"size": 13, "bold": True}),
])
add_slide_number(slide, 21)


# ============================================================
# SLIDE 22: LIMITATIONS & FUTURE WORK
# ============================================================
slide = prs.slides.add_slide(blank_layout)
slide_title(slide, "Limitations & Future Work")
add_multi_text(slide, Inches(0.7), Inches(1.6), Inches(5.5), Inches(5.5), [
    ("Limitations", {"size": 20, "bold": True, "color": RED, "space_after": 10}),
    ("Dataset age (2015)", {"size": 15, "bold": True, "space_after": 2}),
    ("Misses crypto boom, AI wave, COVID effects.", {"size": 13, "color": GRAY, "space_after": 10}),
    ("Limited text data", {"size": 15, "bold": True, "space_after": 2}),
    ("Only category tags, not free-text descriptions.", {"size": 13, "color": GRAY, "space_after": 10}),
    ("Static snapshots", {"size": 15, "bold": True, "space_after": 2}),
    ("3 time slices, not continuous temporal modeling.", {"size": 13, "color": GRAY, "space_after": 10}),
    ("Bipartite only", {"size": 15, "bold": True, "space_after": 2}),
    ("No investor-investor or company-company edges.", {"size": 13, "color": GRAY, "space_after": 10}),
    ("Survivorship bias", {"size": 15, "bold": True, "space_after": 2}),
    ("Only companies visible in Crunchbase.", {"size": 13, "color": GRAY}),
])
add_multi_text(slide, Inches(7.0), Inches(1.6), Inches(5.5), Inches(5.5), [
    ("Future Work", {"size": 20, "bold": True, "color": GREEN, "space_after": 10}),
    ("Dynamic temporal graphs (TGN / TGAT)", {"size": 15, "bold": True, "space_after": 2}),
    ("Model investment timing, not just existence.", {"size": 13, "color": GRAY, "space_after": 10}),
    ("Richer text features", {"size": 15, "bold": True, "space_after": 2}),
    ("Company descriptions, founder bios, press releases.", {"size": 13, "color": GRAY, "space_after": 10}),
    ("Co-investment edges", {"size": 15, "bold": True, "space_after": 2}),
    ("Model investor-investor relationships.", {"size": 13, "color": GRAY, "space_after": 10}),
    ("Node2Vec as XGBoost features", {"size": 15, "bold": True, "space_after": 2}),
    ("Give XGBoost graph structure it can't see today.", {"size": 13, "color": GRAY, "space_after": 10}),
    ("Product deployment", {"size": 15, "bold": True, "space_after": 2}),
    ("REST API for live predictions,", {"size": 13, "color": GRAY}),
    ("investor dashboard, similarity search.", {"size": 13, "color": GRAY}),
])
add_slide_number(slide, 22)


# ============================================================
# SLIDE 23: WHAT EACH MODEL SEES (VISUAL SUMMARY)
# ============================================================
slide = prs.slides.add_slide(blank_layout)
slide_title(slide, "Summary: What Each Model Actually Sees", "Side-by-side comparison")

# XGBoost column
add_text_box(slide, Inches(0.7), Inches(1.5), Inches(3.5), Inches(0.35),
             "XGBoost", font_size=18, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(0.7), Inches(1.9), Inches(3.5), Inches(0.015), BLUE)
add_multi_text(slide, Inches(0.7), Inches(2.1), Inches(3.5), Inches(4.5), [
    ("Sees:", {"size": 14, "bold": True, "space_after": 4}),
    ("One row per company:", {"size": 13}),
    ("  raised=$500K", {"size": 12, "color": GRAY}),
    ("  investors=2", {"size": 12, "color": GRAY}),
    ("  sector=fintech", {"size": 12, "color": GRAY}),
    ("  is_US=yes", {"size": 12, "color": GRAY, "space_after": 8}),
    ("Does NOT see:", {"size": 14, "bold": True, "color": RED, "space_after": 4}),
    ("  Who the investors are", {"size": 12, "color": GRAY}),
    ("  Their track record", {"size": 12, "color": GRAY}),
    ("  Other portfolio companies", {"size": 12, "color": GRAY}),
    ("  Network position", {"size": 12, "color": GRAY}),
])

# GraphSAGE column
add_text_box(slide, Inches(4.8), Inches(1.5), Inches(3.8), Inches(0.35),
             "GraphSAGE", font_size=18, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(4.8), Inches(1.9), Inches(3.8), Inches(0.015), BLUE)
add_multi_text(slide, Inches(4.8), Inches(2.1), Inches(3.8), Inches(4.5), [
    ("Sees everything XGBoost sees, PLUS:", {"size": 14, "bold": True, "space_after": 4}),
    ("  Layer 1 (direct neighbors):", {"size": 13, "bold": True}),
    ("    Investor A: 15 investments,", {"size": 12, "color": GRAY}),
    ("    $5M median, 10 years active", {"size": 12, "color": GRAY}),
    ("    Investor B: 2 investments,", {"size": 12, "color": GRAY}),
    ("    $100K median, 1 year active", {"size": 12, "color": GRAY, "space_after": 4}),
    ("  Layer 2 (2-hop neighbors):", {"size": 13, "bold": True}),
    ("    Investor A's OTHER companies:", {"size": 12, "color": GRAY}),
    ("    12 of 15 raised follow-on!", {"size": 12, "color": GREEN}),
    ("    They're in similar sectors.", {"size": 12, "color": GRAY, "space_after": 4}),
    ("  Learns: 'strong investor syndicate'", {"size": 13, "bold": True, "color": GREEN}),
])

# Ensemble column
add_text_box(slide, Inches(9.2), Inches(1.5), Inches(3.5), Inches(0.35),
             "Ensemble", font_size=18, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(9.2), Inches(1.9), Inches(3.5), Inches(0.015), BLUE)
add_multi_text(slide, Inches(9.2), Inches(2.1), Inches(3.5), Inches(4.5), [
    ("Combines both views:", {"size": 14, "bold": True, "space_after": 6}),
    ("XGBoost:", {"size": 13, "bold": True}),
    ("  'Low funding, few investors'", {"size": 12, "color": RED}),
    ("  P = 0.45 -> NO", {"size": 12, "color": RED, "space_after": 6}),
    ("GraphSAGE:", {"size": 13, "bold": True}),
    ("  'But great investor syndicate!'", {"size": 12, "color": GREEN}),
    ("  P = 0.72 -> YES", {"size": 12, "color": GREEN, "space_after": 6}),
    ("Ensemble:", {"size": 13, "bold": True}),
    ("  P = 0.4*0.45 + 0.6*0.72", {"size": 12, "color": GRAY}),
    ("  P = 0.61 -> YES", {"size": 14, "bold": True, "color": GREEN, "space_after": 6}),
    ("The network context overrides", {"size": 13}),
    ("the weak tabular signal.", {"size": 13}),
])

add_text_box(slide, Inches(0.7), Inches(6.7), Inches(12.0), Inches(0.4),
             "This is why graph structure matters: it provides context that no spreadsheet can capture.",
             font_size=15, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 23)


# ============================================================
# SLIDE 24: KEY TAKEAWAYS
# ============================================================
slide = prs.slides.add_slide(blank_layout)
slide_title(slide, "Key Takeaways")
takeaways = [
    "Graph structure of VC networks carries real predictive signal\nbeyond what tabular features capture (+3pp AUC).",
    "GraphSAGE is the right GNN architecture for bipartite graphs.\nGCN and GAT fail due to structural incompatibilities.",
    "Features dominate (~20pp AUC), but graph is complementary.\nThe ensemble captures both perspectives.",
    "Temporal integrity is essential and validated.\nNo leakage detected; methodology is sound.",
    "Best model catches 3 out of 4 successful startups (75% recall)\nat seed stage. Useful for screening, not decision-making.",
]
for i, text in enumerate(takeaways):
    y = Inches(1.6) + Inches(i * 1.05)
    add_text_box(slide, Inches(0.7), y + Inches(0.03), Inches(0.4), Inches(0.4),
                 str(i + 1) + ".", font_size=16, bold=True, color=BLUE)
    add_text_box(slide, Inches(1.3), y, Inches(11.0), Inches(0.85), text, font_size=14, color=BLACK)
add_rect(slide, Inches(0.7), Inches(6.6), Inches(12.0), Inches(0.02), BLUE)
add_text_box(slide, Inches(0.7), Inches(6.75), Inches(12.0), Inches(0.4),
             "Total improvement:  F1  0.558 -> 0.638  (+14.3%)   |   AUC  0.711 -> 0.745  (+4.8%)",
             font_size=16, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)
add_slide_number(slide, 24)


# ============================================================
# SLIDE 25: THANK YOU
# ============================================================
slide = prs.slides.add_slide(blank_layout)
add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.0),
             "Thank You", font_size=42, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)
add_rect(slide, Inches(5.5), Inches(3.5), Inches(2.3), Inches(0.025), BLUE)
add_text_box(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.6),
             "Questions?", font_size=24, color=GRAY, alignment=PP_ALIGN.CENTER)

# Small network at the bottom
draw_network(slide, Inches(4.5), Inches(4.8), scale=0.7)

add_slide_number(slide, 25)


# Save
output_path = "/Volumes/ThunderDB/ML_Investor_Company_Raise/presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path} ({TOTAL_SLIDES} slides)")
